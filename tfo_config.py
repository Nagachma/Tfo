from pptx.enum.chart import XL_CHART_TYPE , XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.dml.line import LineFormat
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

#GRAPH1 constants cpp sample -bubble chart

median_text="Median CPP for filtered selection - {}"
average_text="Average CPP for filtered selection - {}"

# ---add spPr--- chart border fix.

chartborder_color="000000" #BLACK
chartborder_width=1
spPr_xml = (
    "<c:spPr %s %s>\n"
    "  <a:noFill/>\n"
    "  <a:ln>\n"
    "    <a:solidFill>\n"
    "      <a:srgbClr val=\"DEDEDE\"/>\n"
    "    </a:solidFill>\n"
    "  </a:ln>\n"
    "  <a:effectLst/>\n"
    "</c:spPr>\n" % (nsdecls("c"), nsdecls("a"))
)


bubble_weights=dict()
bubble_weights["Active"]=3
bubble_weights["Completed"]=1
bubble_weights["Planned –Trial scenario"]=2

series_colorsmap_dict=dict()
series_colorsmap_dict["Active"]="6495ED" #blue
series_colorsmap_dict["Completed"]="00C957" #green
series_colorsmap_dict["Planned –Trial scenario"]="CD2626" #red

series_cpp_similartrials_color="FFC125" #goldenrod1


study_status=["Active","Completed","Planned –Trial scenario"]

#GRAPH2 constants cpp similar trials -column chart
cpp_divide_value=1000

fill_transparency=0.75  # sets opacity to 25%

#common constants
data_labels_number_format = '\$00.00\k'
data_labels_font_color="000000" #black
# "FFFFFF" #white

#chart attributes or features.
charttitle_column_dict=dict()
charttitle_column_dict["cpp_sample"]="DEVELOPMENT_UNIT"
charttitle_column_dict["cpp_similartrials"]=None
charttitle_column_dict["regional_allocation"]=None


legendposition_dict=dict()
legendposition_dict["cpp_sample"]=XL_LEGEND_POSITION.TOP
legendposition_dict["cpp_similartrials"]=XL_LEGEND_POSITION.CORNER
legendposition_dict["regional_allocation"]= XL_LEGEND_POSITION.RIGHT


include_in_layout_dict=dict()
include_in_layout_dict["cpp_sample"]=False
include_in_layout_dict["cpp_similartrials"]=True
include_in_layout_dict["regional_allocation"]=False

seriesname_dict=dict()
seriesname_dict["cpp_sample"]=None
seriesname_dict["cpp_similartrials"]="External"
#"Color by cost category External"

xaxis_title_dict=dict()
xaxis_title_dict["cpp_sample"]="# Patients/Subjects"
xaxis_title_dict["cpp_similartrials"]="Trial Code"

xaxis_numberformat_dict=dict()
xaxis_numberformat_dict["cpp_sample"]=None
xaxis_numberformat_dict["cpp_similartrials"]=None

yaxis_title_dict=dict()
yaxis_title_dict["cpp_sample"]="Lifetime CPP"
yaxis_title_dict["cpp_similartrials"]="Lifetime CPP"

yaxis_numberformat_dict=dict()
yaxis_numberformat_dict["cpp_sample"]='\$00\k'
yaxis_numberformat_dict["cpp_similartrials"]='\$00\k'


placeholders_dict=dict()
placeholders_dict["cpp_sample"]=11
placeholders_dict["cpp_similartrials"]=12
placeholders_dict["regional_allocation"]=11
placeholders_dict["cpp_sampletext"]=13

xl_chart_type_dict=dict()
xl_chart_type_dict["cpp_sample"]=XL_CHART_TYPE.BUBBLE
xl_chart_type_dict["cpp_similartrials"]=XL_CHART_TYPE.COLUMN_CLUSTERED
xl_chart_type_dict["regional_allocation"]=XL_CHART_TYPE.PIE

chart_hastitle_dict=dict()
chart_hastitle_dict["cpp_sample"]=True
chart_hastitle_dict["cpp_similartrials"]=False
chart_hastitle_dict["regional_allocation"]=True

fontsize_dict=dict()
fontsize_dict["cpp_sample"]=8
fontsize_dict["cpp_similartrials"]=8
fontsize_dict["regional_allocation"]=8

charttitle_font=8
charttextplaceholder_font=8

datalabels_numberformat_dict=dict()
datalabels_numberformat_dict["cpp_sample"]='\$00.00\k'
datalabels_numberformat_dict["cpp_similartrials"]='\$00.00\k'

legend_dict=dict()
legend_dict["cpp_sample"]=True
legend_dict["cpp_similartrials"]=True
legend_dict["regional_allocation"]=True


chart_title_dict=dict()
chart_title_dict["cpp_sample"]=""
chart_title_dict["cpp_similartrials"]=""
chart_title_dict["regional_allocation"]='Regional  Allocation' 


slide8_id=258
# slide8_template_path=r"services\TFO_slides_graphs\tfo_templates\TFM_Trials_slide8_template_withoutfooter.pptx"
slide8_template_path=r"services\TFO_slides_graphs\tfo_templates\oncology T1 Global Interventional TFO Template 0128.pptx"
    # TFM_Trials_slide8_template_withfooter.pptx"

columns_G_AB=["DRUG_SUBSTANCE_DRUG_PRODUCT",
            "RAW_PACK_LAB_INTERMEDIATES",
            "COMPARATOR_COSTS",
            "INVESTIGATOR_FEES",
            "LAB_ANALYTICAL_SERVICES","SCIENTIFIC_MEETING_COSTS",
            "CRO_COSTS", "PRINTING_CRF_LABELS_BOOKLETS","OTHER_PROJECT_REL_STUDY_COSTS",
            "CONSULTANTS_ACADEMIA_COLLAB","REGISTRATION_AND_LICENSE_FEES",
            "NP4_GRANTS_FUNDING_ARRANG", "PACKAGING","LOGISTICS","DEVICE_DEVELOPMENT_SUPPLY",
            "PROCESS_PACKAGING_DEVELOPMENT","MEDICAL_COMMUNICATIONS","DM_COST",
            "SP_COST","RWS_COST","TM_COST","TMO_COST"
            ]

slide8output_filename="TFM_Trials Forecast Slide8 {} .pptx"


#slide10 piechart constants.
slide10output_filename="TFM_Trials Forecast Slide10 {} .pptx"

slide10_id=256
slide10_template_path=r"services\TFO_slides_graphs\tfo_templates\TFM_Trials_slide10_template.pptx"


allocation_list=["Europe North","AMEA","LATAM",
                 "North America","Europe Central","Europe South East & Baltics","Australia"]

allocation_columnname="Number of Planned Site"

slide10_categories=["EU","LATAM","US","Australia","AMEA"]

slide10_category_columnmap=dict()
slide10_category_columnmap["EU"]=["Europe North","Europe Central","Europe South East & Baltics"]
slide10_category_columnmap["LATAM"] =["LATAM"]
slide10_category_columnmap["US"]=["North America"]
slide10_category_columnmap["Australia"]=["Australia"]
slide10_category_columnmap["AMEA"]=["AMEA"]

slide10_total="Total"
slide10_seriesname='Allocation'
slide10_sheetname='Country_Allocation'

#Embed excel in pptx constants
static_columnnames_list=['A', 'B', 'C', 'D', 'E', 'F', 'G', 'H', 'I', 'J', 'K', 'L', 'M', 'N', 'O', 'P', 'Q', 'R', 'S', 'T', 'U', 'V', 'W',
'X', 'Y', 'Z', 'AA', 'AB', 'AC', 'AD', 'AE', 'AF', 'AG', 'AH', 'AI', 'AJ', 'AK', 'AL', 'AM', 'AN', 'AO', 'AP', 'AQ', 'AR', 'AS', 'AT', 'AU', 'AV', 'AW', 'AX', 'AY', 'AZ', 'BA', 'BB', 'BC', 'BD', 'BE', 'BF', 'BG', 'BH', 'BI', 'BJ', 'BK', 'BL', 'BM', 'BN', 'BO', 'BP', 'BQ', 'BR', 'BS', 'BT', 'BU', 'BV', 'BW', 'BX', 'BY', 'BZ', 'CA', 'CB', 'CC', 'CD', 'CE', 'CF', 'CG', 'CH', 'CI', 'CJ', 'CK', 'CL', 'CM', 'CN', 'CO', 'CP', 'CQ', 'CR', 'CS', 'CT', 'CU', 'CV', 'CW', 'CX', 'CY', 'CZ', 'DA', 'DB', 'DC', 'DD', 'DE', 'DF', 'DG', 'DH', 'DI', 'DJ', 'DK', 'DL', 'DM', 'DN', 'DO', 'DP', 'DQ', 'DR', 'DS', 'DT', 'DU', 'DV', 'DW', 'DX', 'DY', 'DZ', 'EA', 'EB', 'EC', 'ED', 'EE', 'EF', 'EG', 'EH', 'EI', 'EJ', 'EK', 'EL', 'EM', 'EN', 'EO', 'EP', 'EQ', 'ER', 'ES', 'ET', 'EU', 'EV', 'EW', 'EX', 'EY', 'EZ', 'FA', 'FB', 'FC', 'FD', 'FE', 'FF', 'FG', 'FH', 'FI', 'FJ', 'FK', 'FL', 'FM', 'FN', 'FO', 'FP', 'FQ', 'FR', 'FS', 'FT', 'FU', 'FV', 'FW', 'FX', 'FY', 'FZ', 'GA', 'GB', 'GC', 'GD', 'GE', 'GF', 'GG', 'GH', 'GI', 'GJ', 'GK', 'GL', 'GM', 'GN', 'GO', 'GP', 'GQ', 'GR', 'GS', 'GT', 'GU', 'GV', 'GW', 'GX', 'GY', 'GZ', 'HA', 'HB', 'HC', 'HD', 'HE', 'HF', 'HG', 'HH', 'HI', 'HJ', 'HK', 'HL', 'HM', 'HN', 'HO', 'HP', 'HQ', 'HR', 'HS', 'HT', 'HU', 'HV', 'HW', 'HX', 'HY', 'HZ', 'IA', 'IB', 'IC', 'ID', 'IE', 'IF', 'IG', 'IH', 'II', 'IJ', 'IK', 'IL', 'IM', 'IN', 'IO', 'IP', 'IQ', 'IR', 'IS', 'IT', 'IU', 'IV', 'IW', 'IX', 'IY', 'IZ', 'JA', 'JB', 'JC', 'JD', 'JE', 'JF', 'JG', 'JH', 'JI', 'JJ', 'JK', 'JL', 'JM', 'JN', 'JO', 'JP', 'JQ', 'JR', 'JS', 'JT', 'JU', 'JV', 'JW', 'JX', 'JY', 'JZ', 'KA', 'KB', 'KC', 'KD', 'KE', 'KF', 'KG', 'KH', 'KI', 'KJ', 'KK', 'KL', 'KM', 'KN', 'KO', 'KP', 'KQ', 'KR', 'KS', 'KT', 'KU', 'KV', 'KW', 'KX', 'KY', 'KZ', 'LA', 'LB', 'LC', 'LD', 'LE', 'LF', 'LG', 'LH', 'LI', 'LJ', 'LK', 'LL', 'LM', 'LN', 'LO', 'LP', 'LQ', 'LR', 'LS', 'LT', 'LU', 'LV', 'LW', 'LX', 'LY', 'LZ', 'MA', 'MB', 'MC', 'MD', 'ME', 'MF', 'MG', 'MH', 'MI', 'MJ', 'MK', 'ML', 'MM', 'MN', 'MO', 'MP', 'MQ', 'MR', 'MS', 'MT', 'MU', 'MV', 'MW', 'MX', 'MY', 'MZ', 'NA', 'NB', 'NC', 'ND', 'NE', 'NF', 'NG', 'NH', 'NI', 'NJ', 'NK', 'NL', 'NM', 'NN', 'NO', 'NP', 'NQ', 'NR', 'NS', 'NT', 'NU', 'NV', 'NW', 'NX', 'NY', 'NZ', 'OA', 'OB', 'OC', 'OD', 'OE', 'OF', 'OG', 'OH', 'OI', 'OJ', 'OK', 'OL', 'OM', 'ON', 'OO', 'OP', 'OQ', 'OR', 'OS', 'OT', 'OU', 'OV', 'OW', 'OX', 'OY', 'OZ', 'PA', 'PB', 'PC', 'PD', 'PE', 'PF', 'PG', 'PH', 'PI', 'PJ', 'PK', 'PL', 'PM', 'PN', 'PO', 'PP', 'PQ', 'PR', 'PS', 'PT', 'PU', 'PV', 'PW', 'PX', 'PY', 'PZ', 'QA', 'QB', 'QC', 'QD', 'QE', 'QF', 'QG', 'QH', 'QI', 'QJ', 'QK', 'QL', 'QM', 'QN', 'QO', 'QP', 'QQ', 'QR', 'QS', 'QT', 'QU', 'QV', 'QW', 'QX', 'QY', 'QZ', 'RA', 'RB', 'RC', 'RD', 'RE', 'RF', 'RG', 'RH', 'RI', 'RJ', 'RK', 'RL', 'RM', 'RN', 'RO', 'RP', 'RQ', 'RR', 'RS', 'RT', 'RU', 'RV', 'RW', 'RX', 'RY', 'RZ', 'SA', 'SB', 'SC', 'SD', 'SE', 'SF', 'SG', 'SH', 'SI', 'SJ', 'SK', 'SL', 'SM', 'SN', 'SO', 'SP', 'SQ', 'SR', 'SS', 'ST', 'SU', 'SV', 'SW', 'SX', 'SY', 'SZ', 'TA', 'TB', 'TC', 'TD', 'TE', 'TF', 'TG', 'TH', 'TI', 'TJ', 'TK', 'TL', 'TM', 'TN', 'TO', 'TP', 'TQ', 'TR', 'TS', 'TT', 'TU', 'TV', 'TW', 'TX', 'TY', 'TZ', 'UA', 'UB', 'UC', 'UD', 'UE', 'UF', 'UG', 'UH', 'UI', 'UJ', 'UK', 'UL', 'UM', 'UN', 'UO', 'UP', 'UQ', 'UR', 'US', 'UT', 'UU', 'UV', 'UW', 'UX', 'UY', 'UZ', 'VA', 'VB', 'VC', 'VD', 'VE', 'VF', 'VG', 'VH', 'VI', 'VJ', 'VK', 'VL', 'VM', 'VN', 'VO', 'VP', 'VQ', 'VR', 'VS', 'VT', 'VU', 'VV', 'VW', 'VX', 'VY', 'VZ', 'WA', 'WB', 'WC', 'WD', 'WE', 'WF', 'WG', 'WH', 'WI', 'WJ', 'WK', 'WL', 'WM', 'WN', 'WO', 'WP', 'WQ', 'WR', 'WS', 'WT', 'WU', 'WV', 'WW', 'WX', 'WY', 'WZ', 'XA', 'XB', 'XC', 'XD', 'XE', 'XF', 'XG', 'XH', 'XI', 'XJ', 'XK', 'XL', 'XM', 'XN', 'XO', 'XP', 'XQ', 'XR', 'XS', 'XT', 'XU', 'XV', 'XW', 'XX', 'XY', 'XZ', 'YA', 'YB', 'YC', 'YD', 'YE', 'YF', 'YG', 'YH', 'YI', 'YJ', 'YK', 'YL', 'YM', 'YN', 'YO', 'YP', 'YQ', 'YR', 'YS', 'YT', 'YU', 'YV', 'YW', 'YX', 'YY', 'YZ', 'ZA', 'ZB', 'ZC', 'ZD', 'ZE', 'ZF', 'ZG', 'ZH', 'ZI', 'ZJ', 'ZK', 'ZL', 'ZM', 'ZN', 'ZO', 'ZP', 'ZQ', 'ZR', 'ZS', 'ZT', 'ZU', 'ZV', 'ZW', 'ZX', 'ZY', 'ZZ']

data_range="A1:{}"

first_check="Internal Costs"
second_check="Total External cost of Clinical WP ($k)"
cost_summary_columnname="A"



cppsample_oncologyfilters=["STUDY_TYPE_DESC","DEVELOPMENT_UNIT", "PROVIDING_ORGANIZATION",
                            "INDICATION","PHASE_DERIVED"
                        ]

cppsample_ignorefilters=["INDICATION"]


cpp_slidenames=["CPP_ONCOLOGY","CPP"]

cppsimilartrials_oncologyfilters=["DEVELOPMENT_UNIT","PROVIDING_ORGANIZATION","PHASE_DERIVED",
                                    "STUDY_TYPE_DESC","INDICATION"]

cppsimilartrials_oncoctsfilter=["PROJECT_CODE"]

cppsimilartrails_ignorefilters=["INDICATION"]

#TFo regional allocation piechart constants
ra_slidenames=["TfoRegionalAllocation","RegionalAllocation"]
ra_tfosheetname="Summary"
ra_total="Total"
ra_allocation_columnname="Sites"
ra_categorycolumn ="Region"

