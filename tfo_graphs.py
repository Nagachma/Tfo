import pandas as pd
import time
import traceback
import json
import datetime
import os, sys
import math  
from statistics import median
from lxml import etree
from pptx import Presentation
from pptx import Presentation
from pptx.chart.data import CategoryChartData,ChartData,BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE , XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.dml.line import LineFormat
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

#file imports
from services.TFO_slides_graphs import tfo_config
from services.TFO_slides_graphs.tfoslide10_regionalallocation import create_regional_allocation_chartdata
from services.TFO_slides_graphs.tfograph1_cppsample import create_insert_cppsamplechart
from services.TFO_slides_graphs.tfograph2_cppsimilartrail import create_insert_cppsimilartrailchart
from utils.Logger import get_logger

tfo_graphs_logger = get_logger("tfo_graphs_logger")
tfo_graphs_logger.debug("Initialized the tfo_graphs_logger creation")



def create_tfo_pptx_slide8(input_folder_path,output_folder_path,input_filepath,input_json,input_templatepath):
    prs = Presentation(input_templatepath)
    for slide in prs.slides:
        tfo_graphs_logger.debug(slide.slide_id)
        slide_id=slide.slide_id
    for placeholder in slide.placeholders:
        tfo_graphs_logger.debug(placeholder.placeholder_format.idx)
        tfo_graphs_logger.debug(placeholder.name)
    
    slide=prs.slides.get(slide_id)
    if input_json["SlideName"]=="CPP":
        json_cppsamplekey="CPPSample"
        json_cppsimilartrailskey="CPPSimilarTrials"
    elif input_json["SlideName"]=="CPP_ONCOLOGY":
        json_cppsamplekey="CPPSampleOncologyIR"
        json_cppsimilartrailskey="CPPSimilarTrialsOncologyIR"    
    try:
        #create and insert graph1-cpp sample data- bubble chart in slide8 placeholder
        chart_name="cpp_sample"
        json_payload=input_json[json_cppsamplekey]
        tfo_graphs_logger.debug("Input cpp sample json")
        tfo_graphs_logger.debug(json_payload)
        json_payload["SlideName"]=input_json["SlideName"]
        chart_title=input_json[json_cppsamplekey][tfo_config.charttitle_column_dict[chart_name]]
        planned_studycode_json=input_json
        chart_data ,series_number_dict,median_value,yaxis_list,average_value= create_insert_cppsamplechart(input_filepath, json_payload, chart_name,planned_studycode_json)
        slide=create_general_chart(chart_data,series_number_dict,slide,chart_name,chart_title)
    except Exception as e:
        tfo_graphs_logger.debug(str(traceback.format_exc()))
    try:
        #create and insert graph2 - cpp similar trials - column chart in slide8 placeholder
        chart_name="cpp_similartrials"
        json_payload=input_json[json_cppsimilartrailskey]
        tfo_graphs_logger.debug("Input cpp similar trails json")
        tfo_graphs_logger.debug(json_payload)
        json_payload["SlideName"]=input_json["SlideName"]
        chart_title=tfo_config.chart_title_dict[chart_name]
        planned_studycode=input_json[json_cppsamplekey]["STUDY_CODE"]
        planned_cpp=input_json["CPP"]
        tfo_graphs_logger.debug("cpp value passed")
        tfo_graphs_logger.debug(planned_cpp)
        chart_data ,series_number_dict= create_insert_cppsimilartrailchart(input_filepath, json_payload, chart_name,planned_studycode,planned_cpp)
        slide=create_general_chart(chart_data,series_number_dict,slide,chart_name,chart_title)
    except Exception as e:
        tfo_graphs_logger.debug(str(traceback.format_exc()))

    try:
        #insert median and average values to text placeholder.
        sentences_list=[]
        sentences_list.append(tfo_config.median_text.format(round(median_value,1)))
        sentences_list.append(tfo_config.average_text.format(round(average_value,1)))
        placeholder_id=tfo_config.placeholders_dict["cpp_sampletext"]
        slide=insert_text_chart(slide, placeholder_id,sentences_list)
    except Exception as e:
        tfo_graphs_logger.debug(str(traceback.format_exc()))
        
    output_filename=tfo_config.slide8output_filename.format(str(datetime.datetime.now().strftime("%d_%m_%Y-%H_%M_%S")))
    output_file_path=output_folder_path+"/"+output_filename
    prs.save(output_file_path)
    return output_file_path, output_filename,median_value,yaxis_list,average_value

def create_general_chart(chart_data,series_number_dict,slide,chart_name,chart_title):
    '''
    creates bubble chart and column Clustered charts
    '''
    xl_chart_type=tfo_config.xl_chart_type_dict[chart_name]
    # add chart to slide        
    placeholder = slide.placeholders[tfo_config.placeholders_dict[chart_name]]
    chart = placeholder.insert_chart(
        xl_chart_type, chart_data
    ).chart
    
    if chart_name=="cpp_sample":
        series_colorsmap_dict=tfo_config.series_colorsmap_dict
        # Go through every series to modify color
        for series_index, series_name in series_number_dict.items():
            chart.series[series_index].format.fill.solid()
            chart.series[series_index].format.fill.fore_color.rgb = RGBColor.from_string(series_colorsmap_dict[series_name])
            chart.series[series_index].format.fill.transparency = tfo_config.fill_transparency 

        #add border to cppsample chart.
        plotArea = chart._chartSpace.plotArea
        # ---get-or-add spPr---
        spPrs = plotArea.xpath("c:spPr")
        tfo_graphs_logger.debug(len(spPrs))
        if len(spPrs) > 0:
            tfo_graphs_logger.debug("-----len greater then zero-----")
            spPr = spPrs[0]
        else:
            tfo_graphs_logger.debug("--no spPrs tags------")
            # ---add spPr---
            spPr_xml = tfo_config.spPr_xml
            spPr = parse_xml(spPr_xml)
            plotArea.insert_element_before(spPr, "c:chart")

        line = LineFormat(spPr)
        line.color.rgb = RGBColor.from_string(tfo_config.chartborder_color)
        line.width = Pt(tfo_config.chartborder_width)


    if chart_name=="cpp_similartrials":
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = RGBColor.from_string(tfo_config.series_cpp_similartrials_color)
        #Legend to show series name rather then category.
        chart.plots[0].vary_by_categories = False

    chart.has_legend = tfo_config.legend_dict[chart_name]
    if tfo_config.legend_dict[chart_name]==True:
        chart.legend.position = tfo_config.legendposition_dict[chart_name]
        chart.legend.include_in_layout =tfo_config.include_in_layout_dict[chart_name]

    if tfo_config.chart_hastitle_dict[chart_name]==True:
        chart.chart_title.has_text_frame=True
        if isinstance(chart_title,list):
            chart_title=", ".join(chart_title)
        chart.chart_title.text_frame.text=str(chart_title)
        chart.chart_title.text_frame.paragraphs[0].font.bold=True
        chart.chart_title.text_frame.paragraphs[0].font.size=Pt(tfo_config.charttitle_font)
             
    else:
        chart.has_title = False
    
    if chart_name!="regional_allocation":
        #add xaxis title
        chart.category_axis.axis_title.text_frame.text= tfo_config.xaxis_title_dict[chart_name]
        #add yaxis title
        chart.value_axis.axis_title.text_frame.text= tfo_config.yaxis_title_dict[chart_name]
        chart.value_axis.tick_labels.number_format =tfo_config.yaxis_numberformat_dict[chart_name]
        chart.value_axis.has_major_gridlines=False
        chart.value_axis.has_minor_gridlines=False

        #xaxis yaxis title font
        font_size=tfo_config.fontsize_dict[chart_name]
        chart.font.size=Pt(font_size)
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(font_size)
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size =Pt(font_size)

        #data labels format and font    
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(font_size)
        data_labels.font.color.rgb = RGBColor.from_string(tfo_config.data_labels_font_color)
        data_labels.number_format = tfo_config.datalabels_numberformat_dict[chart_name]
        data_labels.position = XL_LABEL_POSITION.CENTER
        
    return slide


def insert_text_chart(slide, placeholder_id,sentences_list):
    #insert text for objectives
    text_placeholder = slide.placeholders[placeholder_id]
    TP=text_placeholder.text_frame
    for sentence in sentences_list:
        paragraph=TP.add_paragraph()
        paragraph.text=sentence    
        paragraph.font.size = Pt(tfo_config.charttextplaceholder_font)
    TP.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    TP.word_wrap = True
    return slide


def create_tfo_pptx_slide10(input_folder_path,output_folder_path,input_filepath,input_json,input_templatepath,slide_name):
    prs = Presentation(input_templatepath)
    for slide in prs.slides:
        tfo_graphs_logger.debug(slide.slide_id)
        slide_id=slide.slide_id
        
    for placeholder in slide.placeholders:
        tfo_graphs_logger.debug(placeholder.placeholder_format.idx)
    
    slide=prs.slides.get(slide_id)

    #create and insert regional allocation- pie chart in slide10 placeholder
    chart_name="regional_allocation"
    chart_data, series_number_dict = create_regional_allocation_chartdata(input_filepath,chart_name,slide_name)
    chart_title=tfo_config.chart_title_dict[chart_name]
    slide=create_general_chart(chart_data,series_number_dict,slide,chart_name,chart_title)

    output_filename=tfo_config.slide10output_filename.format(str(datetime.datetime.now().strftime("%d_%m_%Y-%H_%M_%S")))
    output_file_path=output_folder_path+"/"+output_filename
    prs.save(output_file_path)
    return output_file_path, output_filename




